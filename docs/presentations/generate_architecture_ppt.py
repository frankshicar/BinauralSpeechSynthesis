#!/usr/bin/env python3
"""
Generate PowerPoint presentation for GeoWarpFiLMNet architecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title):
    """Add a content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=None):
    """Add a text box to slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    return textbox

def add_bullet_points(slide, left, top, width, height, points, font_size=14):
    """Add bullet points"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = point
        p.level = 0
        p.font.size = Pt(font_size)
    
    return textbox

def create_architecture_ppt():
    """Create the main presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "GeoWarpFiLMNet v6.3",
                   "Binaural Audio Synthesis with FiLM")
    
    # Slide 2: Overview
    slide = add_content_slide(prs, "模型概覽")
    points = [
        "目標：從單聲道音訊合成雙耳音訊",
        "輸入：mono (B,1,9600) + view (B,7,24)",
        "輸出：y_L, y_R (B,1,9600)",
        "核心技術：",
        "  • Geometric + Neural Warp（時域處理）",
        "  • FiLM ResBlocks（頻域調製）",
        "  • Magnitude/Phase 分離預測",
        "參數量：2,203,718"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 16)
    
    # Slide 3: Complete Architecture
    slide = add_content_slide(prs, "完整架構流程")
    arch_text = """
mono (B,1,9600) + view (B,7,24)
    ↓
┌─────────────────────┐
│  GeometricWarper    │  物理公式計算 ITD
└─────────────────────┘
    ↓ geo_wf (B,2,9600)
┌─────────────────────┐
│ NeuralWarpCorrector │  6層Conv1d修正
└─────────────────────┘
    ↓ y_init (B,2,9600)
┌─────────────────────┐
│       STFT          │  n_fft=1024, hop=256
└─────────────────────┘
    ↓ Y_L_init, Y_R_init (複數)
┌─────────────────────┐
│  提取 Magnitude     │  .abs()
└─────────────────────┘
    ↓ (B,2,513,38)
┌─────────────────────┐
│      Encoder        │  Conv2d: 2→128 channels
└─────────────────────┘
    ↓ (B,128,513,38)
┌─────────────────────┐
│  FiLM ResBlocks ×8  │  根據位置調製特徵
└─────────────────────┘
    ↓ (B,128,513,38)
┌──────────┬──────────┐
│Magnitude │  Phase   │
│   Head   │   Head   │
└──────────┴──────────┘
    ↓           ↓
  mag_pred   phase_res
    ↓           ↓
  重建複數頻譜 Y_pred
    ↓
  iSTFT
    ↓
  y_L, y_R
"""
    add_text_box(slide, 0.5, 1.5, 9, 5.5, arch_text, 11, False, RGBColor(0, 0, 0))
    
    # Slide 4: Warp Components
    slide = add_content_slide(prs, "Warp 組件（串聯）")
    points = [
        "1. GeometricWarper（物理公式，無參數）",
        "   • 計算：geo_wf = -distance / 343 * 48000",
        "   • 輸出：geo_wf (B, 2, 9600)",
        "",
        "2. NeuralWarpCorrector（神經網路，有參數）",
        "   • 輸入：mono + view + geo_wf",
        "   • 6層 Conv1d → delta",
        "   • warpfield = geo_wf + delta",
        "   • MonotoneTimeWarper → y_init",
        "",
        "設計理念：物理公式提供初始值，神經網路學習修正"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 14)
    
    # Slide 5: FiLM Mechanism
    slide = add_content_slide(prs, "FiLM 調製機制")
    film_text = """
核心公式：output = gamma × input + beta

流程：
1. view → TemporalPositionEncoder → pos_feat (B,256,38)

2. 每個 ResBlock：
   pos_feat → film_net (Conv1d 256→128) → params
   params → reshape → gamma, beta (B,64,38)
   gamma, beta → 擴展到 513 頻率 → (B,513,38)
   
3. 調製：
   x = Conv2d(x_in)           # 提取特徵
   x = gamma × x + beta       # FiLM 調製
   x = ReLU(x + residual)     # Residual connection

作用：根據位置動態調整不同頻率的特徵
"""
    add_text_box(slide, 0.5, 1.5, 9, 5.5, film_text, 13, False, RGBColor(0, 0, 0))
    
    # Slide 6: FiLM ResBlock Detail
    slide = add_content_slide(prs, "FiLM ResBlock 詳細結構")
    resblock_text = """
輸入：x_in (B,128,513,38) + pos_feat (B,256,38)

┌─────────────────────────────────────┐
│  residual = x_in  (保存原始輸入)     │
│                                     │
│  ┌──────────────────────────────┐  │
│  │  FiLM Parameter Net          │  │
│  │  pos_feat → Conv1d → γ, β    │  │
│  └──────────────────────────────┘  │
│              ↓                      │
│  ┌──────────────────────────────┐  │
│  │  ResBlock 主分支              │  │
│  │  x = Conv2d(x_in)             │  │
│  │  x = BatchNorm + ReLU         │  │
│  │  x = Conv2d(x)                │  │
│  │  x = BatchNorm                │  │
│  │  x = gamma × x + beta  (FiLM) │  │
│  └──────────────────────────────┘  │
│              ↓                      │
│  x_out = ReLU(x + residual)        │
└─────────────────────────────────────┘

輸出：x_out (B,128,513,38)
"""
    add_text_box(slide, 0.5, 1.5, 9, 5.5, resblock_text, 12, False, RGBColor(0, 0, 0))
    
    # Slide 7: Decoder Heads
    slide = add_content_slide(prs, "Decoder：兩個輸出頭")
    points = [
        "輸入：FiLM ResBlocks 輸出 (B,128,513,38)",
        "",
        "MagnitudeHead：",
        "  • Conv2d 128→64→2",
        "  • Softplus 激活（確保 > 0）",
        "  • 輸出：mag_L, mag_R (B,513,38)",
        "",
        "PhaseHead：",
        "  • Conv2d 128→64→2",
        "  • Tanh × π 激活（範圍 [-π, π]）",
        "  • 輸出：phase_res (B,2,513,38)",
        "",
        "重建：",
        "  • phase_pred = phase_init + phase_res",
        "  • Y_pred = mag_pred × exp(j × phase_pred)",
        "  • y = iSTFT(Y_pred)"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 13)
    
    # Slide 8: Key Design Choices
    slide = add_content_slide(prs, "關鍵設計選擇")
    points = [
        "1. 為什麼只用 Magnitude 做輸入？",
        "   • Magnitude 是平滑的、非負的，更適合 CNN",
        "   • Phase 可以從 magnitude 特徵推斷",
        "",
        "2. 為什麼 Phase 用 Residual？",
        "   • phase_init 已經接近正確（來自 warp）",
        "   • 只需學習小的修正，更容易訓練",
        "",
        "3. 為什麼分兩個 Head？",
        "   • Magnitude 和 Phase 特性不同",
        "   • 分開學習更容易，可以用不同的激活函數",
        "",
        "4. 為什麼用 FiLM？",
        "   • 根據位置動態調整特徵",
        "   • 每層都可以做位置相關的調整"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 12)
    
    # Slide 9: v6.3 Improvements
    slide = add_content_slide(prs, "v6.3 改進")
    points = [
        "1. Log-scale 頻帶分配",
        "   • v6：線性分配，低頻只有 4/64 參數",
        "   • v6.3：對數分配，低頻有 ~20/64 參數",
        "   • 效果：低頻控制更精細",
        "",
        "2. NeuralWarpCorrector 加深",
        "   • v6：4 層 Conv1d",
        "   • v6.3：6 層 Conv1d",
        "   • 效果：更大感受野，更好的 ITD 預測",
        "",
        "3. Low-freq sin²(diff/2) IPD Loss",
        "   • 只計算前 32 個頻率 bin",
        "   • sin²(diff/2) 比 cosine 有更好的梯度",
        "   • 效果：低頻相位預測更準確"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 13)
    
    # Slide 10: Data Flow Summary
    slide = add_content_slide(prs, "數據流總結")
    flow_text = """
時域 → 頻域 → 特徵 → 調製 → 預測 → 重建

mono (時域)
  → Warp → y_init (時域，雙耳)
  → STFT → Y_init (頻域，複數)
  → .abs() → magnitude (頻域，實數)
  → Encoder → 特徵 (128 channels)
  → FiLM → 調製後的特徵
  → Heads → magnitude + phase (預測)
  → 重建 → Y_pred (頻域，複數)
  → iSTFT → y (時域，雙耳)

關鍵：
• Warp 處理 ITD（時域）
• FiLM 處理 ILD/IPD（頻域）
• Magnitude 特徵 → 推斷 Phase
"""
    add_text_box(slide, 0.5, 1.5, 9, 5.5, flow_text, 14, False, RGBColor(0, 0, 0))
    
    # Slide 11: L/R Channel Flow
    slide = add_content_slide(prs, "左右聲道處理流程")
    points = [
        "分開階段：",
        "  • Warp 輸出：y_init (B,2,9600) - 通道維度分 L/R",
        "  • STFT 後：Y_L_init, Y_R_init - 兩個獨立 tensor",
        "  • 最終輸出：y_L, y_R - 兩個獨立 tensor",
        "",
        "合併階段：",
        "  • Encoder 輸入：stack([mag_L, mag_R]) → (B,2,513,38)",
        "  • Encoder 輸出：(B,128,513,38) - L/R 混合成 128 特徵",
        "  • FiLM 處理：(B,128,513,38) - 處理混合特徵",
        "",
        "再分開階段：",
        "  • Magnitude Head：(B,2,513,38) - 通道 0=L, 1=R",
        "  • Phase Head：(B,2,513,38) - 通道 0=L, 1=R",
        "",
        "設計理念：合併處理學習 L/R 關係，最後分開輸出"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 12)
    
    # Slide 12: Comparison with DPATFNet
    slide = add_content_slide(prs, "與 DPATFNet 的差異")
    comp_text = """
項目                  DPATFNet              GeoWarpFiLMNet v6.3
─────────────────────────────────────────────────────────
STFT hop_length      64 (高時間解析度)      256 (高頻率解析度)
時間幀數              150 frames            38 frames
STFT 後處理          保留複數頻譜           只用 magnitude
頻域處理方式          Upsampling            固定解析度 FiLM
Position 融合        MPF (淺層融合)         FiLM (每層調製)
架構風格              生成式                 編碼-解碼
Doppler 建模         DPAB (dual attention)  FiLM (implicit)
輸出方式              Conv1+Tanh            重建複數頻譜+iSTFT

性能：
• Phase-L2: 0.717 (兩者相同)
• 你的模型更簡單，計算量更小
• DPATFNet 時間解析度更高，更適合捕捉 Doppler effect
"""
    add_text_box(slide, 0.3, 1.5, 9.4, 5.5, comp_text, 11, False, RGBColor(0, 0, 0))
    
    # Slide 13: Training Details
    slide = add_content_slide(prs, "訓練細節")
    points = [
        "Loss Function：",
        "  • Stage 1: L2 + Low-freq sin²(diff/2) IPD",
        "  • Stage 2: L2 + Phase + IPD",
        "",
        "STFT 參數：",
        "  • n_fft: 1024",
        "  • hop_length: 256",
        "  • win_length: 1024",
        "",
        "模型參數：",
        "  • Encoder: Conv2d 2→128 channels",
        "  • FiLM ResBlocks: 8 層，dilation [1,2,4,8,1,2,4,8]",
        "  • NeuralWarpCorrector: 6 層 Conv1d",
        "  • 總參數: 2,203,718"
    ]
    add_bullet_points(slide, 0.5, 1.5, 9, 5, points, 14)
    
    # Save presentation
    output_path = "/home/sbplab/frank/BinauralSpeechSynthesis/GeoWarpFiLMNet_Architecture.pptx"
    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")
    return output_path

if __name__ == "__main__":
    create_architecture_ppt()
